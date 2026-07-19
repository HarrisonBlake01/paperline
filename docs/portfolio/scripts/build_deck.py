#!/usr/bin/env python3
"""Build Paperline recruiter deck (10 core slides + 1 appendix)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

A = "/tmp/pl_assets/"
OUT = "/Users/openclaw-server/.openclaw/workspace/paperline/app/docs/portfolio/Paperline_Recruiter_Case_Study.pptx"

# ---- palette ----
BG      = RGBColor(0x0B, 0x0D, 0x12)
CARD    = RGBColor(0x15, 0x1A, 0x24)
CARD2   = RGBColor(0x1B, 0x21, 0x2E)
BORDER  = RGBColor(0x2A, 0x31, 0x40)
WHITE   = RGBColor(0xF5, 0xF7, 0xFA)
GRAY    = RGBColor(0x9A, 0xA4, 0xB5)
BLUE    = RGBColor(0x5B, 0x8D, 0xEF)
GREEN   = RGBColor(0x3D, 0xC9, 0x8E)
AMBER   = RGBColor(0xF0, 0xB4, 0x29)
FONT    = "Helvetica Neue"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def txt(slide, x, y, w, h, runs, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=None, wrap=True):
    """runs: str, or list of paragraphs; each paragraph is str or list of (text, dict) runs."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = space_after
        if isinstance(para, str):
            para = [(para, {})]
        for text, style in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = style.get("font", FONT)
            f.size = Pt(style.get("size", size))
            f.bold = style.get("bold", bold)
            f.color.rgb = style.get("color", color)
            if style.get("link"):
                r.hyperlink.address = style["link"]
    return box

def card(slide, x, y, w, h, fill=CARD, line=BORDER, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.paragraphs[0].text = ""
    return sh

def chip(slide, x, y, w, h, text, color, size=13, fill=None, bold=True):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sh.adjustments[0] = 0.5
    except Exception:
        pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill if fill else CARD2
    sh.line.color.rgb = color; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sh

def pic(slide, path, x, y, w=None, h=None, desc=None, border=True):
    p = slide.shapes.add_picture(path, x, y, w, h)
    if desc:
        p._element._nvXxPr.cNvPr.set('descr', desc)
    if border:
        p.line.color.rgb = BORDER; p.line.width = Pt(1)
    return p

def fit_w(path, target_h):
    im = Image.open(path); return Emu(int(target_h * im.width / im.height))

def fit_h(path, target_w):
    im = Image.open(path); return Emu(int(target_w * im.height / im.width))

def page_num(slide, n):
    txt(slide, Inches(12.75), Inches(7.08), Inches(0.45), Inches(0.3), str(n),
        size=12, color=GRAY, align=PP_ALIGN.RIGHT)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def kicker(slide, text, color=BLUE, x=Inches(0.6), y=Inches(0.42)):
    txt(slide, x, y, Inches(8), Inches(0.3), text.upper(), size=13, color=color, bold=True)

def headline(slide, text, y=Inches(0.72), size=32, w=Inches(12.1)):
    txt(slide, Inches(0.6), y, w, Inches(0.85), text, size=size, color=WHITE, bold=True)

# ================= SLIDE 1 — TITLE =================
s = new_slide()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.09), SH)
accent.fill.solid(); accent.fill.fore_color.rgb = BLUE; accent.line.fill.background(); accent.shadow.inherit = False
txt(s, Inches(0.75), Inches(1.15), Inches(6.0), Inches(1.1), "Paperline", size=54, bold=True)
txt(s, Inches(0.75), Inches(2.25), Inches(5.9), Inches(1.3),
    "Cited document intelligence with human-controlled operations", size=24, color=BLUE, bold=True, line_spacing=1.05)
txt(s, Inches(0.75), Inches(3.55), Inches(5.7), Inches(0.6),
    "Full-stack Applied AI portfolio case study · Harrison Olvera", size=18, color=GRAY)
chip(s, Inches(0.75), Inches(4.25), Inches(4.35), Inches(0.44), "Synthetic recruiter demo · Stripe test mode", AMBER, size=13)
pic(s, A+"hero-jobcard.png", Inches(7.15), Inches(0.85), w=Inches(5.45), h=fit_h(A+"hero-jobcard.png", Inches(5.45)),
    desc="Paperline Ops Agent operations job card showing two synthetic documents and the Upload, Extract, Approve, Operate workflow trace")
txt(s, Inches(0.75), Inches(6.75), Inches(11.8), Inches(0.4),
    "Next.js · TypeScript · Clerk · Supabase/pgvector · OpenAI · Stripe", size=14, color=GRAY)
notes(s, "Paperline is my flagship portfolio project. It addresses the gap between an AI model reading a document and a team safely acting on the result. The interface you see here is the live Ops Agent demo route: two synthetic documents grouped into one operations job, moving through an upload, extract, approve, operate workflow. Everything consequential stays behind a human approval boundary, and billing runs in Stripe test mode. Over the next few minutes I'll walk through the product problem, the working architecture, and where I drew the line between working code, synthetic demo data, and planned productionization.")

# ================= SLIDE 2 — PROBLEM =================
s = new_slide()
kicker(s, "The product problem")
headline(s, "Document AI needs evidence, structure, and control—not just plausible answers.", size=30)
cards = [
    ("Evidence", "Can the reviewer trace a field or answer back to the source page?",
     "Quoted snippets + page-level citations", BLUE),
    ("Structure", "Can the output enter a repeatable workflow instead of remaining chat text?",
     "Schema-guided, runtime-validated extraction", GREEN),
    ("Control", "Can the system prepare actions without silently spending or changing external systems?",
     "Human approval before spend or provisioning", AMBER),
]
cw, ch, gap = Inches(3.93), Inches(4.1), Inches(0.25)
x0, y0 = Inches(0.6), Inches(2.15)
for i, (t, q, a, c) in enumerate(cards):
    x = x0 + i * (cw + gap)
    card(s, x, y0, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, Inches(0.07), ch)
    bar.fill.solid(); bar.fill.fore_color.rgb = c; bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, x+Inches(0.35), y0+Inches(0.35), cw-Inches(0.7), Inches(0.5), t, size=24, bold=True, color=c)
    txt(s, x+Inches(0.35), y0+Inches(1.05), cw-Inches(0.7), Inches(1.7), q, size=18, color=WHITE, line_spacing=1.15)
    txt(s, x+Inches(0.35), y0+Inches(3.0), cw-Inches(0.7), Inches(0.95),
        [[("Paperline's answer  ", {"size":13, "color":GRAY, "bold":True})], [(a, {"size":15, "color":c})]],
        line_spacing=1.1)
page_num(s, 2)
notes(s, "I framed this as an engineering problem, not a demo problem. Teams get plausible answers from document AI, but three things are usually missing. First, evidence: a reviewer should be able to trace any extracted field to a quoted snippet and a page number. Second, structure: output has to be schema-validated so it can enter a repeatable workflow instead of staying chat text. Third, control: the system can prepare billing or vendor actions, but it must not silently spend money or change external systems. Underneath that sit mixed file types, tenant boundaries, uncertainty signals, and approval ownership—those constraints drove the architecture.")

# ================= SLIDE 3 — WORKFLOW =================
s = new_slide()
kicker(s, "End-to-end workflow")
headline(s, "From upload to approval-ready operations")
chip(s, Inches(8.35), Inches(0.42), Inches(4.45), Inches(0.44), "Synthetic recruiter demo — not a live customer workflow", AMBER, size=11)
steps = [
    ("01", "Upload", "2 synthetic documents", "MSA + vendor invoice"),
    ("02", "Extract", "4 cited fields", "values · quotes · pages"),
    ("03", "Approve", "3 recommended actions", "owners + approval states"),
    ("04", "Operate", "Stripe test-mode preview", "$6.00 · 20 pages · paused"),
]
cw, ch, gap = Inches(2.82), Inches(2.5), Inches(0.32)
x0, y0 = Inches(0.6), Inches(1.9)
for i, (num, t, big, sub) in enumerate(steps):
    x = x0 + i * (cw + gap)
    card(s, x, y0, cw, ch)
    txt(s, x+Inches(0.3), y0+Inches(0.28), cw-Inches(0.6), Inches(0.35), num, size=14, color=BLUE, bold=True)
    txt(s, x+Inches(0.3), y0+Inches(0.66), cw-Inches(0.6), Inches(0.55), t, size=26, bold=True)
    txt(s, x+Inches(0.3), y0+Inches(1.38), cw-Inches(0.6), Inches(0.55), big, size=17, color=GREEN if i<3 else AMBER, bold=True, line_spacing=1.05)
    txt(s, x+Inches(0.3), y0+Inches(2.0), cw-Inches(0.6), Inches(0.4), sub, size=13, color=GRAY)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x+cw+Inches(0.02), y0+Inches(1.05), Inches(0.28), Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = BLUE; ar.line.fill.background(); ar.shadow.inherit = False
pic(s, A+"hero-metrics.png", Inches(0.6), Inches(4.85), w=Inches(12.13), h=fit_h(A+"hero-metrics.png", Inches(12.13)),
    desc="Ops Agent metric strip: 2 documents, 4 cited fields, 3 approval-ready agent actions, $6.00 Stripe test mode billing preview")
txt(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
    "Live demo route: paperline-xi.vercel.app/ops-agent — one reviewer-friendly surface for the whole product story.",
    size=14, color=GRAY)
page_num(s, 3)
notes(s, "The Ops Agent route compresses the product story into a single reviewer-friendly surface. Two synthetic documents—a master services agreement and a vendor invoice—flow through four stages: upload, extract, approve, operate. Extraction produces four cited fields with quoted snippets and page references. The agent proposes three actions, each with an owner and an approval state. And the operate step shows a six-dollar Stripe test-mode usage preview for twenty processed pages. I want to be explicit: this is a synthetic recruiter demo, not a live customer workflow—the point is showing the exact shape of the workflow the working product paths support.")

# ================= SLIDE 4 — TRUSTWORTHY EXTRACTION =================
s = new_slide()
kicker(s, "Trustworthy extraction")
headline(s, "Facts with receipts")
img_h = Inches(4.7)
pic(s, A+"cited-approvals.png", Inches(0.6), Inches(1.62), w=fit_w(A+"cited-approvals.png", img_h), h=img_h,
    desc="Cited extraction cards with confidence labels, quoted snippets and page references, next to recommended actions awaiting human approval")
callouts = [
    ("Page-level references", "Every field cites its source document and page.", BLUE),
    ("Quoted source snippets", "The exact supporting text ships with the value.", BLUE),
    ("Confidence = review signal", "Labels prioritize review—not a guarantee.", GREEN),
    ("Human review flag", "Auto-renewal clause is explicitly flagged for a person.", AMBER),
]
cy = Inches(1.62)
for t, d, c in callouts:
    card(s, Inches(9.5), cy, Inches(3.25), Inches(0.98), fill=CARD2)
    txt(s, Inches(9.7), cy+Inches(0.11), Inches(2.95), Inches(0.35), t, size=14, bold=True, color=c, wrap=False)
    txt(s, Inches(9.7), cy+Inches(0.46), Inches(2.9), Inches(0.5), d, size=12, color=GRAY, line_spacing=1.0)
    cy += Inches(1.13)
chip(s, Inches(0.6), Inches(6.55), Inches(7.6), Inches(0.44),
     "Displayed documents, values, and confidence labels are synthetic fixtures", AMBER, size=13)
page_num(s, 4)
notes(s, "Evidence is part of the product interface, not something hidden in logs or a developer console. Each extraction card shows the value, a quoted snippet of the source text, the document name, and the page reference—so a reviewer can verify the fact without reopening the file. Confidence labels are presented as a review signal, and the auto-renewal clause carries an explicit human-review flag because that's exactly the kind of term a person should confirm. One important caveat: the documents, values, and confidence labels on this screen are synthetic fixtures built to illustrate the workflow—not a live extraction run.")

# ================= SLIDE 5 — DETERMINISTIC CONTROLS =================
s = new_slide()
kicker(s, "Deterministic controls around AI")
headline(s, "The model extracts; the application controls.")
lw = Inches(6.0)
card(s, Inches(0.6), Inches(1.75), lw, Inches(4.9))
txt(s, Inches(0.95), Inches(2.05), Inches(5.3), Inches(0.45), "MODEL-ASSISTED", size=15, bold=True, color=BLUE)
model_items = ["Document classification", "Schema-guided extraction", "OCR-assisted processing of scans and images", "Grounded answer generation with citations"]
iy = Inches(2.65)
for it in model_items:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), iy+Inches(0.09), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = BLUE; dot.line.fill.background(); dot.shadow.inherit = False
    txt(s, Inches(1.25), iy, Inches(5.1), Inches(0.5), it, size=18)
    iy += Inches(0.62)
txt(s, Inches(0.95), Inches(5.6), Inches(5.3), Inches(0.9),
    "Probabilistic output is treated as a proposal:\ninspectable, cited, and validated before use.", size=14, color=GRAY, line_spacing=1.15)
card(s, Inches(6.85), Inches(1.75), Inches(5.88), Inches(4.9))
txt(s, Inches(7.2), Inches(2.05), Inches(5.2), Inches(0.45), "DETERMINISTIC APPLICATION CONTROLS", size=15, bold=True, color=GREEN)
det_items = ["Authentication + workspace authorization", "Runtime schema validation of extraction output", "Private storage and persisted records", "Usage and billing rules", "Approval ownership and state"]
iy = Inches(2.65)
for it in det_items:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), iy+Inches(0.09), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background(); dot.shadow.inherit = False
    txt(s, Inches(7.5), iy, Inches(5.0), Inches(0.5), it, size=18)
    iy += Inches(0.62)
txt(s, Inches(7.2), Inches(5.85), Inches(5.2), Inches(0.6),
    "Model judgment never owns identity, money, tenancy, or state.", size=14, color=GRAY)
txt(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.35),
    "Boundary map verified against repository source: src/lib/ai/*, src/lib/auth/workspace.ts, src/app/api/billing/*, supabase/migrations/*",
    size=12, color=GRAY)
page_num(s, 5)
notes(s, "This is the design principle I care most about: knowing where probabilistic model behavior should stop. The model handles what models are good at—classification, schema-guided extraction, OCR-assisted processing, and grounded answers. Everything consequential is deterministic application code: Clerk authentication and workspace authorization, runtime schema validation of every extraction shape, private storage and persistence, usage and billing rules, and approval ownership and state. The model proposes; the application decides. That split is verifiable in the repository—the AI code lives in src/lib/ai, and the control paths live in the auth, billing, and migration layers.")

# ================= SLIDE 6 — ARCHITECTURE =================
s = new_slide()
kicker(s, "System architecture")
headline(s, "Full-stack document intelligence architecture")
# legend
lg = [("Implemented", GREEN), ("Synthetic demo", AMBER), ("Planned", GRAY)]
lx = Inches(8.35)
for name, c in lg:
    sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, Inches(0.52), Inches(0.18), Inches(0.18))
    sq.fill.solid(); sq.fill.fore_color.rgb = c; sq.line.fill.background(); sq.shadow.inherit = False
    txt(s, lx+Inches(0.26), Inches(0.47), Inches(1.5), Inches(0.3), name, size=12, color=GRAY, wrap=False)
    lx += Inches(1.75) if name != "Implemented" else Inches(1.55)

def node(x, y, w, h, title, sub, color, dashed=False, tsize=14, ssize=10.5, sub_y=0.42):
    sh = card(s, x, y, w, h, fill=CARD2, line=color, radius=0.12)
    if dashed:
        ln = sh.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    txt(s, x+Inches(0.12), y+Inches(0.08), w-Inches(0.24), Inches(0.32), title, size=tsize, bold=True, color=WHITE)
    if sub:
        txt(s, x+Inches(0.12), y+Inches(sub_y), w-Inches(0.24), h-Inches(sub_y+0.06), sub, size=ssize, color=GRAY, line_spacing=1.0)
    return sh

def arrow(x1, y1, x2, y2, color=BORDER):
    ln = s.shapes.add_connector(2, x1, y1, x2, y2)
    ln.line.color.rgb = color; ln.line.width = Pt(1.5)
    le = ln.line._get_or_add_ln()
    he = le.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    le.append(he)

ROW1 = Inches(1.55); NH = Inches(0.95)
node(Inches(0.6),  ROW1, Inches(1.75), NH, "Signed-in user", "Clerk identity", BLUE)
node(Inches(2.75), ROW1, Inches(2.15), NH, "Next.js routes", "authenticated App Router", GREEN)
node(Inches(5.3),  ROW1, Inches(2.15), NH, "Workspace checks", "membership + role", GREEN)
node(Inches(7.85), ROW1, Inches(2.3),  NH, "Private storage + DB", "Supabase · RLS-aware", GREEN)
node(Inches(10.55),ROW1, Inches(2.2),  NH, "Parsing / OCR", "PDF · DOCX · scans · images", GREEN)
arrow(Inches(2.35), ROW1+NH/2, Inches(2.75), ROW1+NH/2)
arrow(Inches(4.9),  ROW1+NH/2, Inches(5.3),  ROW1+NH/2)
arrow(Inches(7.45), ROW1+NH/2, Inches(7.85), ROW1+NH/2)
arrow(Inches(10.15),ROW1+NH/2, Inches(10.55),ROW1+NH/2)
ROW2 = Inches(3.05)
node(Inches(0.6),  ROW2, Inches(2.6), NH, "Chunks + embeddings", "page-aware · pgvector", GREEN)
node(Inches(3.6),  ROW2, Inches(2.9), NH, "Extraction + cited chat", "schema-validated · page citations", GREEN)
node(Inches(6.9),  ROW2, Inches(2.7), NH, "Workflow records", "persisted · synchronous today", GREEN)
node(Inches(10.0), ROW2, Inches(2.75), NH, "Billing", "Stripe Checkout/Portal + signed webhook", GREEN)
arrow(Inches(3.2),  ROW2+NH/2, Inches(3.6), ROW2+NH/2)
arrow(Inches(6.5),  ROW2+NH/2, Inches(6.9), ROW2+NH/2)
arrow(Inches(9.6),  ROW2+NH/2, Inches(10.0),ROW2+NH/2)
ROW3 = Inches(4.85)
card(s, Inches(0.6), Inches(4.55), Inches(6.1), Inches(2.15), fill=CARD)
txt(s, Inches(0.85), Inches(4.68), Inches(5.6), Inches(0.3), "OPS AGENT DEMO SURFACE — SYNTHETIC FIXTURES", size=12, bold=True, color=AMBER)
node(Inches(0.85), Inches(5.1), Inches(1.7), Inches(0.9), "Fixture docs", "MSA + invoice", AMBER, dashed=True, tsize=12.5)
node(Inches(2.7),  Inches(5.1), Inches(1.85), Inches(0.9), "Actions + approvals", "owners · states", AMBER, dashed=True, tsize=11, sub_y=0.55)
node(Inches(4.7),  Inches(5.1), Inches(1.8), Inches(0.9), "$6 test preview", "no real charge", AMBER, dashed=True, tsize=11, sub_y=0.55)
txt(s, Inches(0.85), Inches(6.1), Inches(5.6), Inches(0.5),
    "Illustrates the workflow shape — never mistaken for live integration.", size=11, color=GRAY)
card(s, Inches(7.0), Inches(4.55), Inches(5.73), Inches(2.15), fill=CARD)
txt(s, Inches(7.25), Inches(4.68), Inches(5.2), Inches(0.3), "PLANNED PRODUCTIONIZATION — NOT YET WIRED", size=12, bold=True, color=GRAY)
node(Inches(7.25), Inches(5.1), Inches(1.6), Inches(0.9), "Inngest", "durable execution", GRAY, dashed=True, tsize=12.5)
node(Inches(9.0),  Inches(5.1), Inches(1.75), Inches(0.9), "Sentry/PostHog", "observability", GRAY, dashed=True, tsize=11.5)
node(Inches(10.9), Inches(5.1), Inches(1.7), Inches(0.9), "NemoClaw /\nOpenShell", "secure runtime", GRAY, dashed=True, tsize=10.5, sub_y=0.6)
txt(s, Inches(7.25), Inches(6.1), Inches(5.2), Inches(0.5),
    "Explicit next steps — documented, not claimed as complete.", size=11, color=GRAY)
txt(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.35),
    "Service-role DB access bypasses RLS — route-level workspace filters remain an explicit, tested control.",
    size=13, color=GRAY)
page_num(s, 6)
notes(s, "The top two rows are the implemented product path, all present in the repository: a signed-in user hits authenticated Next.js routes, workspace membership and role checks resolve tenancy, documents land in private Supabase storage, and parsing or OCR feeds page-aware chunks and embeddings into pgvector. From there, schema-validated extraction and cited chat produce reviewable output, workflow records persist—execution is synchronous today—and Stripe Checkout, Portal, and a signature-verified webhook handle billing. The amber box is the synthetic Ops Agent demo surface, and the gray box is planned productionization—kept visually separate so nobody mistakes them for live integration. One trust-boundary detail I always mention: server-side service-role access bypasses RLS, so route-level workspace filters are a critical control I treat as code to review and test, not an assumption.")

# ================= SLIDE 7 — EVALUATION =================
s = new_slide()
kicker(s, "Evaluation, not vibes")
headline(s, "A deterministic extraction regression harness")
chip(s, Inches(0.6), Inches(1.6), Inches(8.9), Inches(0.46),
     "Sample-prediction baseline used to validate the scorer — not live-model or production accuracy", AMBER, size=13)
metrics = [
    ("3", "synthetic document cases", GRAY),
    ("20", "labeled fields", GRAY),
    ("40.00%", "exact accuracy", BLUE),
    ("75.00%", "normalized accuracy", BLUE),
    ("94.12%", "presence F1", GREEN),
    ("85.71%", "list-item F1", GREEN),
]
cw, ch, gap = Inches(1.96), Inches(1.55), Inches(0.075)
x0, y0 = Inches(0.6), Inches(2.35)
for i, (v, l, c) in enumerate(metrics):
    x = x0 + i * (cw + gap)
    card(s, x, y0, cw, ch)
    txt(s, x, y0+Inches(0.22), cw, Inches(0.6), v, size=28, bold=True, color=c, align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.1), y0+Inches(0.92), cw-Inches(0.2), Inches(0.55), l, size=13, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=0.95)
# exact vs normalized comparison
card(s, Inches(0.6), Inches(4.3), Inches(7.6), Inches(2.3))
txt(s, Inches(0.9), Inches(4.5), Inches(7.0), Inches(0.35), "WHY EXACT MATCHING UNDERSTATES QUALITY", size=13, bold=True, color=BLUE)
bary = Inches(5.05)
for label, frac, val in [("Exact", 0.40, "40.00%  (8/20)"), ("Normalized", 0.75, "75.00%  (15/20)")]:
    txt(s, Inches(0.9), bary, Inches(1.55), Inches(0.3), label, size=15, color=WHITE)
    track = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), bary+Inches(0.03), Inches(3.9), Inches(0.28))
    track.adjustments[0] = 0.5; track.fill.solid(); track.fill.fore_color.rgb = CARD2
    track.line.color.rgb = BORDER; track.line.width = Pt(0.75); track.shadow.inherit = False
    fill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), bary+Inches(0.03), Emu(int(Inches(3.9) * frac)), Inches(0.28))
    fill.adjustments[0] = 0.5; fill.fill.solid(); fill.fill.fore_color.rgb = BLUE
    fill.line.fill.background(); fill.shadow.inherit = False
    txt(s, Inches(6.55), bary, Inches(1.7), Inches(0.3), val, size=14, color=GRAY)
    bary += Inches(0.62)
txt(s, Inches(0.9), Inches(6.26), Inches(7.0), Inches(0.3),
    "Same predictions, scored two ways — normalization isolates formatting noise.", size=12.5, color=GRAY)
card(s, Inches(8.5), Inches(4.3), Inches(4.23), Inches(2.3))
txt(s, Inches(8.8), Inches(4.5), Inches(3.7), Inches(0.35), "TYPE-AWARE NORMALIZATION", size=13, bold=True, color=GREEN)
txt(s, Inches(8.8), Inches(4.95), Inches(3.7), Inches(1.6),
    ["Casing and whitespace", "Currency and numbers", "Dates and booleans", "Order-insensitive lists"],
    size=15, color=WHITE, line_spacing=1.25)
page_num(s, 7)
notes(s, "I didn't want extraction quality judged only by screenshots, so the repository includes a small deterministic regression harness: three synthetic document cases and twenty labeled fields. The numbers here come from deliberately imperfect sample predictions whose job is to validate the scorer itself—they are not production accuracy or live-model quality. The interesting result is the gap: forty percent exact accuracy becomes seventy-five percent once type-aware normalization handles casing, currency, dates, booleans, and list order. That gap is exactly why naive exact matching understates extraction quality. Presence F1 is tracked separately because knowing when a field is genuinely absent—versus missed—deserves its own metric. The harness gives future model and prompt changes a deterministic contract.")

# ================= SLIDE 8 — APPROVALS / BILLING / SECURE RUNTIME =================
s = new_slide()
kicker(s, "Approvals, billing, and secure-runtime direction")
headline(s, "Agentic value without uncontrolled action", size=30)
img_h = Inches(3.05)
img_w = fit_w(A+"billing-trace-security.png", img_h)
pic(s, A+"billing-trace-security.png", Emu(int((SW - img_w) / 2)), Inches(1.55), w=img_w, h=img_h,
    desc="Stripe test-mode business operation card with waiting-for-human-approval notice, Hermes operator trace, and NVIDIA secure runtime path panel")
cols = [
    ("DEMONSTRATED ON THE SITE", GREEN,
     ["Visible action owners and states", "Stripe test-mode usage preview", "External spend/provisioning paused"]),
    ("IMPLEMENTED IN THE REPOSITORY", BLUE,
     ["Stripe Checkout / Portal routes", "Signature-verified subscription webhook", "Persisted workflow records"]),
    ("VALIDATED PATH — NOT YET INTEGRATED", GRAY,
     ["NVIDIA NemoClaw/OpenShell policy-governed Hermes sandbox", "Controlled network access", "Managed credential/inference handling"]),
]
cw, ch, gap = Inches(3.94), Inches(1.85), Inches(0.24)
x0, y0 = Inches(0.6), Inches(4.8)
for i, (t, c, items) in enumerate(cols):
    x = x0 + i * (cw + gap)
    card(s, x, y0, cw, ch)
    txt(s, x+Inches(0.25), y0+Inches(0.14), cw-Inches(0.5), Inches(0.35), t, size=12.5, bold=True, color=c)
    txt(s, x+Inches(0.25), y0+Inches(0.52), cw-Inches(0.5), Inches(1.25),
        [[("· "+it, {"size":13, "color":WHITE})] for it in items], line_spacing=1.12)
txt(s, Inches(0.6), Inches(6.78), Inches(12.1), Inches(0.55),
    [[("NVIDIA NemoClaw/OpenShell is a validated secure-runtime direction, not a completed Paperline integration.  ", {"size":14, "color":AMBER, "bold":True}),
      ("Source: NVIDIA/NemoClaw and NVIDIA/OpenShell official GitHub documentation, verified 2026-07-14.", {"size":12, "color":GRAY})]])
page_num(s, 8)
notes(s, "This slide separates three layers of honesty. What's demonstrated on the site: every proposed action shows an owner and an approval state, the billing card is explicitly Stripe test mode, and external spend and provisioning are paused pending a human decision. What's implemented elsewhere in the repository: real Stripe Checkout and Portal routes, a signature-verified subscription webhook, and persisted workflow records. And the third column is a validated architecture direction that is not yet integrated: running the Hermes operator inside an NVIDIA NemoClaw/OpenShell sandbox with policy-governed network access and managed credential and inference handling. What makes that direction credible is that official NVIDIA documentation explicitly supports Hermes through NemoClaw. But I want to be precise: Paperline has not wired that runtime yet, and NVIDIA does not secure the current deployment.")

# ================= SLIDE 9 — NEXT STEPS =================
s = new_slide()
kicker(s, "Engineering tradeoffs and next steps")
headline(s, "What I would productionize next")
cards9 = [
    ("1", "Durable execution", "Move synchronous workflow execution to Inngest with idempotency, retries, resume, and dead-letter handling."),
    ("2", "Persisted approvals", "Add explicit approval policies and state transitions from proposed through completed/failed."),
    ("3", "Evaluation expansion", "Capture live predictions against a larger versioned, privacy-safe benchmark; add retrieval/citation evaluation."),
    ("4", "Observability + isolation", "Design safe Sentry/PostHog events and validate the NemoClaw/OpenShell path without logging document content."),
]
cw, ch = Inches(6.0), Inches(2.15)
poss = [(Inches(0.6), Inches(1.9)), (Inches(6.85), Inches(1.9)), (Inches(0.6), Inches(4.3)), (Inches(6.85), Inches(4.3))]
for (n, t, d), (x, y) in zip(cards9, poss):
    card(s, x, y, cw, ch)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.3), y+Inches(0.32), Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = BLUE; circ.line.fill.background(); circ.shadow.inherit = False
    tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    txt(s, x+Inches(1.0), y+Inches(0.34), cw-Inches(1.3), Inches(0.5), t, size=21, bold=True)
    txt(s, x+Inches(0.35), y+Inches(1.0), cw-Inches(0.7), Inches(1.0), d, size=15.5, color=GRAY, line_spacing=1.15)
txt(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.35),
    "Roadmap items — designed and documented, deliberately not claimed as complete.", size=13, color=GRAY)
page_num(s, 9)
notes(s, "These are the four things I would productionize first, in priority order. Durable execution: workflows currently run synchronously in one request, so moving to Inngest gets idempotency, retries, resume, and dead-letter handling. Persisted approvals: the approval boundary should become explicit policies and state transitions from proposed through completed or failed. Evaluation expansion: I want live predictions captured against a larger versioned, privacy-safe benchmark, plus retrieval and citation evaluation. And observability with isolation: safe Sentry and PostHog events designed so document content never leaks into logs, alongside validating the NemoClaw/OpenShell path. None of these are done—being explicit about limitations is part of how I demonstrate production judgment.")

# ================= SLIDE 10 — CLOSE =================
s = new_slide()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.09), SH)
accent.fill.solid(); accent.fill.fore_color.rgb = BLUE; accent.line.fill.background(); accent.shadow.inherit = False
kicker(s, "Recruiter close", x=Inches(0.75))
headline(s, "What Paperline demonstrates", y=Inches(0.75))
proofs = [
    "Full-stack AI product engineering",
    "Cited retrieval and schema-guided extraction",
    "Multi-tenant auth/data-boundary thinking",
    "Human-in-the-loop agent design",
    "Evaluation, billing, and productionization judgment",
]
py = Inches(1.95)
for pr in proofs:
    chk = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), py+Inches(0.06), Inches(0.3), Inches(0.3))
    chk.fill.solid(); chk.fill.fore_color.rgb = GREEN; chk.line.fill.background(); chk.shadow.inherit = False
    tfc = chk.text_frame; pp = tfc.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = "✓"; rr.font.size = Pt(13); rr.font.bold = True; rr.font.color.rgb = BG; rr.font.name = FONT
    txt(s, Inches(1.35), py, Inches(6.2), Inches(0.45), pr, size=20)
    py += Inches(0.72)
txt(s, Inches(0.8), Inches(5.9), Inches(6.5), Inches(0.45),
    "Harrison Olvera · Applied AI / Agentic AI / LLM Engineering", size=18, bold=True, color=WHITE)
card(s, Inches(7.6), Inches(1.95), Inches(5.1), Inches(2.55))
txt(s, Inches(7.95), Inches(2.2), Inches(4.4), Inches(0.35), "SEE IT YOURSELF", size=14, bold=True, color=BLUE)
links = [
    ("Live Ops Agent", "https://paperline-xi.vercel.app/ops-agent"),
    ("GitHub", "https://github.com/HarrisonBlake01/paperline"),
]
ly = Inches(2.7)
for name, url in links:
    txt(s, Inches(7.95), ly, Inches(4.5), Inches(0.35), name, size=16, bold=True, color=WHITE)
    disp = url.replace("https://", "")
    if len(disp) > 46: disp = disp[:44] + "…"
    txt(s, Inches(7.95), ly+Inches(0.34), Inches(4.5), Inches(0.32),
        [[(disp, {"size":13, "color":RGBColor(0x8A, 0xB4, 0xF8), "link":url})]])
    ly += Inches(0.88)
txt(s, Inches(0.8), Inches(6.8), Inches(11.9), Inches(0.35),
    "All demo data synthetic · Stripe test mode · no customer, revenue, or compliance claims", size=13, color=GRAY)
page_num(s, 10)
notes(s, "To close: Paperline demonstrates full-stack AI product engineering—not a wrapper around a chat call. Cited retrieval and schema-guided extraction make output verifiable. Workspace-scoped auth and data design show multi-tenant boundary thinking. The approval boundary shows human-in-the-loop agent design. And the eval harness, billing paths, and explicit roadmap show judgment about what production actually requires. The live Ops Agent page and the GitHub repository are linked here. I'm happy to go deeper on any layer—the extraction pipeline, the tenancy model, the eval design, or the secure-runtime direction.")

# ================= SLIDE 11 — APPENDIX =================
s = new_slide()
kicker(s, "Appendix", color=GRAY)
headline(s, "Claims, sources, and boundaries")
rows = [
    ("Implemented repository paths", "src/lib/pipeline · src/lib/parsing · src/lib/ai/{ocr,extract,chat} · src/lib/auth/workspace.ts · src/app/api/workflows · src/app/api/billing/{checkout,portal} · src/app/api/webhooks/stripe · supabase/migrations/0001_init.sql, 0002_rls.sql", GREEN),
    ("Synthetic demo boundary", "Ops Agent route uses static fixtures from src/lib/ops-agent-demo.ts: Northstar Supply Co., document names, values, confidence labels, action states, operator trace, and the $6.00 / 20-page Stripe test-mode preview. No live extraction, payment, or external action.", AMBER),
    ("Evaluation caveat", "Metrics are a deliberately imperfect sample-prediction baseline validating the deterministic scorer (3 cases, 20 fields). Not production, live-model, or customer accuracy.", AMBER),
    ("NVIDIA sources", "github.com/NVIDIA/NemoClaw · github.com/NVIDIA/OpenShell · docs.nvidia.com/nemoclaw · docs.nvidia.com/openshell — NemoClaw lists Hermes as a supported agent. Verified 2026-07-14.", BLUE),
    ("Compliance", "Paperline does not claim HIPAA, SOC 2, or legal certification. No customer, revenue, or business-outcome claims.", GRAY),
]
y = Inches(1.7)
for t, d, c in rows:
    card(s, Inches(0.6), y, Inches(12.13), Inches(0.94), fill=CARD)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.06), Inches(0.94))
    bar.fill.solid(); bar.fill.fore_color.rgb = c; bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, Inches(0.9), y+Inches(0.1), Inches(2.75), Inches(0.75), t, size=15, bold=True, color=c, line_spacing=1.0)
    txt(s, Inches(3.8), y+Inches(0.09), Inches(8.75), Inches(0.8), d, size=11.5, color=GRAY, line_spacing=1.03)
    y += Inches(1.06)
txt(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3),
    "Claims ledger: docs/portfolio/presentation-claims-audit.md · Repository gates passed 2026-07-14: test:templates, test:extraction-eval, test:demo, lint, build",
    size=11, color=GRAY)
page_num(s, 11)
notes(s, "This appendix is the claims ledger behind the deck. Every implemented capability maps to a repository path. The synthetic demo boundary lists exactly which values are fixtures. The evaluation numbers are a scorer-validation baseline, not production accuracy. The NVIDIA direction cites official GitHub and docs sources verified on July 14, 2026, and Paperline makes no HIPAA, SOC 2, or legal certification claims. The full audit lives in the repository alongside the deck.")

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
